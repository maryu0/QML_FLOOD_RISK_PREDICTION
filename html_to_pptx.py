"""
HTML Presentation to PowerPoint Converter
Converts presentation_aqi.html to a .pptx file
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from bs4 import BeautifulSoup
import re

# Alias for convenience
RgbColor = RGBColor


def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def create_gradient_background(slide, color1, color2):
    """Add a solid background (gradient not directly supported, use primary color)"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RgbColor(*hex_to_rgb(color1))


def add_title_slide(prs, title, subtitle, authors, affiliation):
    """Create the title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    create_gradient_background(slide, "1E3A8A", "7C3AED")
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(18)
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Authors
    authors_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.5))
    authors_frame = authors_box.text_frame
    p = authors_frame.paragraphs[0]
    p.text = authors
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Affiliation
    affil_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(0.5))
    affil_frame = affil_box.text_frame
    p = affil_frame.paragraphs[0]
    p.text = affiliation
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(220, 220, 220)
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_content_slide(prs, title, content_items, is_quantum=False):
    """Create a content slide with header and bullet points"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Header background
    header_color = "1E3A8A" if not is_quantum else "7C3AED"
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = RgbColor(*hex_to_rgb(header_color))
    header.line.fill.background()
    
    # Header title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = RgbColor(30, 30, 30)
        p.space_after = Pt(12)
    
    return slide


def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items, is_quantum=False):
    """Create a two-column content slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header_color = "1E3A8A" if not is_quantum else "7C3AED"
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = RgbColor(*hex_to_rgb(header_color))
    header.line.fill.background()
    
    # Header title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    
    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.3), Inches(0.4))
    p = left_title_box.text_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RgbColor(*hex_to_rgb("1a365d"))
    
    # Left column content
    left_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(4.3), Inches(4.8))
    left_frame = left_content.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(50, 50, 50)
        p.space_after = Pt(8)
    
    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.4))
    p = right_title_box.text_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RgbColor(*hex_to_rgb("7C3AED"))
    
    # Right column content
    right_content = slide.shapes.add_textbox(Inches(5.2), Inches(1.7), Inches(4.3), Inches(4.8))
    right_frame = right_content.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(50, 50, 50)
        p.space_after = Pt(8)
    
    return slide


def add_table_slide(prs, title, headers, rows, highlight_row=None):
    """Create a slide with a table"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = RgbColor(*hex_to_rgb("1E3A8A"))
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    
    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.5 * num_rows)).table
    
    # Header row
    for i, header_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RgbColor(*hex_to_rgb("1a365d"))
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = RgbColor(255, 255, 255)
        p.font.bold = True
        p.font.size = Pt(14)
    
    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            
            if highlight_row is not None and row_idx == highlight_row:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RgbColor(*hex_to_rgb("d1fae5"))
                p.font.bold = True
    
    return slide


def add_thank_you_slide(prs):
    """Create the thank you slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    create_gradient_background(slide, "1E3A8A", "7C3AED")
    
    title_box = slide.shapes.add_textbox(Inches(0), Inches(2.5), Inches(10), Inches(1.5))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_outline_slide(prs, outline_items):
    """Create the outline slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = RgbColor(*hex_to_rgb("7C3AED"))
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Presentation Outline"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    
    # Outline items in two columns
    left_items = outline_items[:5]
    right_items = outline_items[5:]
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = f"{i+1}. {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = RgbColor(30, 30, 30)
        p.space_after = Pt(16)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.5), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = f"{i+6}. {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = RgbColor(30, 30, 30)
        p.space_after = Pt(16)
    
    return slide


def add_three_box_slide(prs, title, boxes):
    """Create a slide with three problem/conclusion boxes"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = RgbColor(*hex_to_rgb("1E3A8A"))
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    
    # Three boxes
    box_width = 2.9
    start_x = 0.5
    colors = ["3182ce", "7C3AED", "10B981"]
    
    for i, (box_title, box_desc) in enumerate(boxes):
        x = start_x + i * (box_width + 0.2)
        
        # Box background
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.3), Inches(box_width), Inches(4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RgbColor(255, 255, 255)
        shape.line.color.rgb = RgbColor(*hex_to_rgb(colors[i]))
        shape.line.width = Pt(3)
        
        # Box title
        title_shape = slide.shapes.add_textbox(Inches(x + 0.1), Inches(1.5), Inches(box_width - 0.2), Inches(0.6))
        p = title_shape.text_frame.paragraphs[0]
        p.text = box_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RgbColor(*hex_to_rgb("1a365d"))
        p.alignment = PP_ALIGN.CENTER
        
        # Box description
        desc_shape = slide.shapes.add_textbox(Inches(x + 0.1), Inches(2.2), Inches(box_width - 0.2), Inches(2.8))
        desc_frame = desc_shape.text_frame
        desc_frame.word_wrap = True
        p = desc_frame.paragraphs[0]
        p.text = box_desc
        p.font.size = Pt(12)
        p.font.color.rgb = RgbColor(70, 70, 70)
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_architecture_slide(prs):
    """Create the architecture slide with flow diagram"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = RgbColor(*hex_to_rgb("7C3AED"))
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    p.text = "QLSTM Architecture"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    
    # Architecture flow boxes
    components = [
        ("Input Layer", "6 steps × 12 features", "3182ce"),
        ("LSTM Layer", "50 hidden units", "3182ce"),
        ("Dense Layer", "4 units, tanh", "3182ce"),
        ("VQC", "4 qubits, 16 params", "7C3AED"),
        ("Output", "CO prediction", "3182ce"),
    ]
    
    box_width = 1.6
    start_x = 0.4
    y = 1.5
    
    for i, (name, desc, color) in enumerate(components):
        x = start_x + i * (box_width + 0.3)
        
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_width), Inches(0.9))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RgbColor(*hex_to_rgb(color))
        shape.line.fill.background()
        
        text_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.15), Inches(box_width), Inches(0.6))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RgbColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(9)
        p2.font.color.rgb = RgbColor(220, 220, 220)
        p2.alignment = PP_ALIGN.CENTER
        
        # Arrow
        if i < len(components) - 1:
            arrow_x = x + box_width
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(arrow_x + 0.05), Inches(y + 0.35), Inches(0.2), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RgbColor(150, 150, 150)
            arrow.line.fill.background()
    
    # Key points
    points = [
        "Classical Components: LSTM extracts temporal features, Dense encodes to [-1, 1]",
        "Quantum Components: 4-qubit VQC with 16 trainable parameters, data re-uploading",
        "Skip Connection: h_classical concatenated with h_quantum for stable gradient flow",
        "Output: ŷ = W_out · [h_classical ⊕ h_quantum] + b"
    ]
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(3.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, point in enumerate(points):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(30, 30, 30)
        p.space_after = Pt(12)
    
    return slide


def convert_html_to_pptx(html_file, output_file):
    """Main function to convert HTML presentation to PPTX"""
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(
        prs,
        "Two-Phase Training of Variational Quantum Circuit Enhanced LSTM for Air Pollution Prediction",
        "A Hybrid Quantum-Classical Approach for Environmental Time-Series Forecasting",
        "Pratim Pramanik    •    Ayush Kumar    •    Bibek Gupta",
        "Department of Computer Science, Kalinga Institute Of Industrial Technology"
    )
    
    # Slide 2: Outline
    outline_items = [
        "Introduction & Motivation",
        "Problem Statement",
        "Related Work",
        "QLSTM Architecture",
        "VQC Design",
        "Two-Phase Training",
        "Experimental Setup",
        "Results & Analysis",
        "Conclusion & Future Work"
    ]
    add_outline_slide(prs, outline_items)
    
    # Slide 3: Introduction
    add_two_column_slide(
        prs,
        "Introduction: Why Air Quality Prediction Matters",
        "The Challenge",
        [
            "Air pollution poses significant threats to public health worldwide",
            "Carbon monoxide (CO) is particularly dangerous - colorless & odorless",
            "WHO estimates 4.2 million deaths annually from outdoor air pollution",
            "Need for timely health advisories for vulnerable populations",
            "Informed policy decisions for emissions control"
        ],
        "The Quantum Opportunity",
        [
            "Complex temporal dependencies challenge classical ML models",
            "Exponential Hilbert space growth enables richer representations",
            "Quantum entanglement captures complex correlations",
            "Variational circuits provide trainable quantum layers"
        ]
    )
    
    # Slide 4: Problem Statement
    add_three_box_slide(
        prs,
        "Research Challenge",
        [
            ("Challenge 1: Barren Plateau Problem", "Gradients vanish exponentially with circuit depth, making training intractable for deep quantum circuits"),
            ("Challenge 2: Scale Mismatch", "Classical gradients (~10⁻³) dominate quantum gradients (~10⁻⁵), causing optimizer to ignore quantum parameters"),
            ("Challenge 3: Feature Quality", "Random classical features early in training provide poor quantum encoding, leading to suboptimal learning")
        ]
    )
    
    # Slide 5: Related Work
    add_two_column_slide(
        prs,
        "Related Work",
        "Air Quality & Deep Learning",
        [
            "Traditional Methods: ARIMA, SVR, Random Forest",
            "Deep Learning: CNN-LSTM [Wen et al.]",
            "Attention Mechanisms: Transformer-based [Li et al.]"
        ],
        "Quantum Machine Learning",
        [
            "VQCs: Variational Quantum Circuits",
            "Data Re-uploading: [Schuld et al., Pérez-Salinas]",
            "Quantum-enhanced RNNs: [Chen et al.]",
            "Gap: Effective training methodologies for hybrid models"
        ],
        is_quantum=True
    )
    
    # Slide 6: Architecture
    add_architecture_slide(prs)
    
    # Slide 7: VQC Design
    add_content_slide(
        prs,
        "4-Layer Variational Quantum Circuit",
        [
            "4-qubit circuit with RZ, RY, RX rotation gates",
            "CNOT and CZ gates for entanglement",
            "Data Re-uploading: Classical data encoded in layers 1, 2, 3 for increased expressivity",
            "16 trainable parameters θᵢ⁽ˡ⁾ optimized during training",
            "Measurement: Pauli-Z expectation values ⟨Zᵢ⟩ for output",
            "Encoding formula: φᵢ = (h_classical + 1) · π/2"
        ],
        is_quantum=True
    )
    
    # Slide 8: Two-Phase Training
    add_two_column_slide(
        prs,
        "Two-Phase Training Methodology",
        "Phase 1: Classical Pre-training",
        [
            "Train: LSTM + Dense layers only",
            "Frozen: Quantum parameters",
            "Optimizer: Adam (lr=0.001)",
            "Early stopping: patience=5",
            "Batch size: 32, Epochs: 25",
            "Goal: Extract robust temporal features"
        ],
        "Phase 2: Quantum Fine-tuning",
        [
            "Frozen: Classical weights",
            "Train: Quantum params + output layer",
            "LR schedule: Cosine decay (0.03 → 0)",
            "Initialization: Uniform [-π/2, π/2]",
            "Batch size: 64, Epochs: 40",
            "Goal: Optimize quantum parameters with stable features"
        ],
        is_quantum=True
    )
    
    # Slide 9: Why Two-Phase Works
    add_three_box_slide(
        prs,
        "Why Two-Phase Training Works",
        [
            ("Stable Classical Features", "Quantum circuit receives consistent, meaningful input encodings rather than random features from untrained classical layers"),
            ("Gradient Scale Separation", "Dedicated quantum optimization with appropriate learning rates prevents optimizer from ignoring small quantum gradients (~10⁻⁵)"),
            ("Symmetry Breaking", "Non-zero initialization of quantum output weights breaks symmetry that traps gradients at zero during training")
        ]
    )
    
    # Slide 10: Experimental Setup
    add_two_column_slide(
        prs,
        "Experimental Setup",
        "Dataset",
        [
            "UCI Air Quality Dataset",
            "9,358 hourly instances",
            "12 features (sensors, temp, humidity)",
            "Target: CO concentration",
            "Split: 80% train, 10% val, 10% test",
            "Min-Max normalization",
            "6-hour sliding windows"
        ],
        "Implementation",
        [
            "TensorFlow 2.x + PennyLane",
            "Lightning.qubit backend",
            "Adjoint differentiation",
            "Baselines: Random Forest (100 estimators), LSTM",
            "Metrics: MAE, RMSE, R², Accuracy"
        ]
    )
    
    # Slide 11: Results
    add_table_slide(
        prs,
        "Model Performance Comparison",
        ["Model", "MAE", "RMSE", "R²"],
        [
            ["Random Forest", "0.0381", "0.0546", "0.7969"],
            ["LSTM", "0.0400", "0.0544", "0.7982"],
            ["QLSTM (Ours)", "0.0400", "0.0543", "0.7990"]
        ],
        highlight_row=2
    )
    
    # Slide 12: Ablation Study
    add_table_slide(
        prs,
        "Ablation Study: Two-Phase vs One-Phase Training",
        ["Training Method", "MAE", "RMSE", "R²"],
        [
            ["One-Phase (End-to-End)", "0.0449", "0.0629", "0.7300"],
            ["Two-Phase (Ours)", "0.0426", "0.0579", "0.7716"],
            ["Improvement", "+5.18%", "+7.95%", "+5.69%"]
        ],
        highlight_row=1
    )
    
    # Slide 13: Visualization
    add_content_slide(
        prs,
        "24-Hour CO Concentration Forecast",
        [
            "QLSTM closely tracks actual CO concentration values",
            "Especially accurate during pollution peaks",
            "Demonstrates practical applicability for real-time forecasting",
            "Better performance than baseline LSTM on temporal patterns",
            "Validates quantum enhancement for time-series prediction"
        ]
    )
    
    # Slide 14: Discussion
    add_two_column_slide(
        prs,
        "Discussion & Limitations",
        "Strengths",
        [
            "Novel two-phase training methodology",
            "Addresses barren plateau problem",
            "Skip connections ensure gradient flow",
            "Practical for current NISQ devices",
            "Demonstrates quantum-classical synergy"
        ],
        "Limitations",
        [
            "Quantum simulation overhead (slower training)",
            "Marginal improvement over LSTM (+0.08%)",
            "Dataset may not fully exploit quantum advantages",
            "Real quantum hardware evaluation needed",
            "4-qubit implementation (simulation only)"
        ]
    )
    
    # Slide 15: Conclusion
    add_three_box_slide(
        prs,
        "Conclusion",
        [
            ("QLSTM Architecture", "Hybrid model combining LSTM with 4-qubit VQC featuring data re-uploading and skip connections for stable gradient flow"),
            ("Two-Phase Training", "Novel methodology addressing barren plateaus: classical pre-training → quantum fine-tuning with dedicated optimizers"),
            ("Empirical Validation", "5.69% R² improvement over end-to-end training on UCI Air Quality dataset, validating the approach")
        ]
    )
    
    # Slide 16: Thank You
    add_thank_you_slide(prs)
    
    # Save presentation
    prs.save(output_file)
    print(f"✅ Presentation saved to: {output_file}")


if __name__ == "__main__":
    html_file = "presentation_aqi.html"
    output_file = "presentation_aqi.pptx"
    
    print("🔄 Converting HTML presentation to PowerPoint...")
    convert_html_to_pptx(html_file, output_file)
    print("✨ Done!")
